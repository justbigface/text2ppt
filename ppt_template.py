from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import math
import textwrap

# 页面&布局参数
PAGE_WIDTH = Inches(13.33)
PAGE_HEIGHT = Inches(7.5)
TOP_MARGIN = Inches(1.5)
BOTTOM_MARGIN = Inches(1.0)  # 底边留白
CARD_MIN_HEIGHT = Inches(1.2)
CARD_MAX_HEIGHT = Inches(2.2)
CARD_SPACING = Inches(0.5)
DESC_BOX_WIDTH_INCH = 8.2

# 字号参数
FONT_SIZE_LARGE = Pt(16)
FONT_SIZE_MID = Pt(14)
FONT_SIZE_SMALL = Pt(12)

# 单行最大字符数（估算）
LINE_CHAR_COUNT = 34

def estimate_card_height_and_font(desc):
    lines = desc.count('\n') + 1
    length = len(desc)
    # 按字数估算行数
    est_lines = max(lines, math.ceil(length / LINE_CHAR_COUNT))
    # 超长内容则减小字号
    if est_lines <= 2:
        return CARD_MIN_HEIGHT, FONT_SIZE_LARGE
    elif est_lines <= 4:
        return CARD_MIN_HEIGHT + Inches(0.3), FONT_SIZE_MID
    elif est_lines <= 8:
        return CARD_MAX_HEIGHT, FONT_SIZE_SMALL
    else:
        # 超长内容需要分卡片
        return CARD_MAX_HEIGHT, FONT_SIZE_SMALL

def split_desc_to_chunks(desc, max_lines=8):
    # 超过8行时切分
    wrapper = textwrap.TextWrapper(width=LINE_CHAR_COUNT)
    lines = wrapper.wrap(desc)
    chunks = []
    for i in range(0, len(lines), max_lines):
        chunk = "\n".join(lines[i:i+max_lines])
        chunks.append(chunk)
    return chunks

def add_card_item(slide, top, num, subtitle, desc, card_height, font_size):
    left = Inches(1)
    width = Inches(10.5)
    height = card_height
    # 1. 卡片
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(180, 180, 180)

    # 2. 左侧圆
    circle_left = left - Inches(0.65)
    circle_top = top + Inches(0.25)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, circle_top, Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(242, 241, 247)
    circle.line.color.rgb = RGBColor(90, 70, 160)
    circle.line.width = Pt(3)
    # 圆内编号
    num_frame = circle.text_frame
    p = num_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(80, 70, 160)
    p.alignment = PP_ALIGN.CENTER

    # 3. 副标题
    sub_box = slide.shapes.add_textbox(left + Inches(1.1), top + Inches(0.2), Inches(8.5), Inches(0.5))
    sub_frame = sub_box.text_frame
    sub_frame.text = subtitle
    p = sub_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(90, 70, 160)
    p.alignment = PP_ALIGN.LEFT

    # 4. 描述
    desc_box = slide.shapes.add_textbox(left + Inches(1.1), top + Inches(0.7), DESC_BOX_WIDTH_INCH, card_height - Inches(0.8))
    desc_frame = desc_box.text_frame
    desc_frame.text = desc
    p = desc_frame.paragraphs[0]
    p.font.size = font_size
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.alignment = PP_ALIGN.LEFT
