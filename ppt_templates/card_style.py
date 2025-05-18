from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import math
import textwrap

# 页面&布局参数
PAGE_WIDTH = Inches(13.33)
PAGE_HEIGHT = Inches(7.5)
TOP_MARGIN = Inches(1.5)
BOTTOM_MARGIN = Inches(1.0)  # 底边留白
CARD_MIN_HEIGHT = Inches(1.2)
CARD_MAX_HEIGHT = Inches(2.2)
CARD_SPACING = Inches(0.5)
DESC_BOX_WIDTH_INCH = 8.2

# 字号参数
FONT_SIZE_LARGE = Pt(16)
FONT_SIZE_MID = Pt(14)
FONT_SIZE_SMALL = Pt(12)

# 单行最大字符数（估算）
LINE_CHAR_COUNT = 34

def estimate_card_height_and_font(desc):
    lines = desc.count('\n') + 1
    length = len(desc)
    # 按字数估算行数
    est_lines = max(lines, math.ceil(length / LINE_CHAR_COUNT))
    # 超长内容则减小字号
    if est_lines <= 2:
        return CARD_MIN_HEIGHT, FONT_SIZE_LARGE
    elif est_lines <= 4:
        return CARD_MIN_HEIGHT + Inches(0.3), FONT_SIZE_MID
    elif est_lines <= 8:
        return CARD_MAX_HEIGHT, FONT_SIZE_SMALL
    else:
        # 超长内容需要分卡片
        return CARD_MAX_HEIGHT, FONT_SIZE_SMALL

def split_desc_to_chunks(desc, max_lines=8):
    # 超过8行时切分
    wrapper = textwrap.TextWrapper(width=LINE_CHAR_COUNT)
    lines = wrapper.wrap(desc)
    chunks = []
    for i in range(0, len(lines), max_lines):
        chunk = "\n".join(lines[i:i+max_lines])
        chunks.append(chunk)
    return chunks

def add_card_item(slide, top, num, subtitle, desc, card_height, font_size):
    left = Inches(1)
    width = Inches(10.5)
    height = card_height
    # 1. 卡片
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(180, 180, 180)

    # 2. 左侧圆
    circle_left = left - Inches(0.65)
    circle_top = top + Inches(0.25)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, circle_top, Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(242, 241, 247)
    circle.line.color.rgb = RGBColor(90, 70, 160)
    circle.line.width = Pt(3)
    # 圆内编号
    num_frame = circle.text_frame
    p = num_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(80, 70, 160)
    p.alignment = PP_ALIGN.CENTER

    # 3. 副标题
    sub_box = slide.shapes.add_textbox(left + Inches(1.1), top + Inches(0.2), Inches(8.5), Inches(0.5))
    sub_frame = sub_box.text_frame
    sub_frame.text = subtitle
    p = sub_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(90, 70, 160)
    p.alignment = PP_ALIGN.LEFT

    # 4. 描述
    desc_box = slide.shapes.add_textbox(left + Inches(1.1), top + Inches(0.7), DESC_BOX_WIDTH_INCH, card_height - Inches(0.8))
    desc_frame = desc_box.text_frame
    desc_frame.text = desc
    p = desc_frame.paragraphs[0]
    p.font.size = font_size
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.alignment = PP_ALIGN.LEFT

def create_card_ppt(title, items, output_path="output.pptx"):
    prs = Presentation()
    prs.slide_width = PAGE_WIDTH
    prs.slide_height = PAGE_HEIGHT

    new_items = []
    # 先处理超长内容：自动切分成多卡片
    for item in items:
        desc = item.get("desc", "")
        subtitle = item.get("subtitle", "")
        num = item.get("num", "")
        # 拆分desc
        desc_chunks = split_desc_to_chunks(desc)
        for i, chunk in enumerate(desc_chunks):
            _num = f"{num}-{i+1}" if len(desc_chunks) > 1 else num
            new_items.append({
                "num": _num,
                "subtitle": subtitle if i == 0 else "",  # 仅第1张显示副标题
                "desc": chunk
            })
    items = new_items

    # 计算每页能容纳的卡片数
    avail_height = PAGE_HEIGHT - TOP_MARGIN - BOTTOM_MARGIN
    idx = 0
    total = len(items)
    slide = None
    row = 0
    used_height = 0
    while idx < total:
        if slide is None or used_height + CARD_MIN_HEIGHT > avail_height:
            # 新建幻灯片
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            used_height = 0
            row = 0
            # 只在第一页加标题
            if len(prs.slides) == 1:
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = title
                p = title_frame.paragraphs[0]
                p.font.size = Pt(40)
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT

        item = items[idx]
        # 估算高度和字号
        card_height, font_size = estimate_card_height_and_font(item["desc"])
        # 防止最后一条超页高，强制拆分
        if used_height + card_height > avail_height:
            slide = None
            continue
        top = TOP_MARGIN + used_height
        add_card_item(
            slide,
            top=top,
            num=item.get("num", idx + 1),
            subtitle=item.get("subtitle", ""),
            desc=item.get("desc", ""),
            card_height=card_height,
            font_size=font_size
        )
        used_height += card_height + CARD_SPACING
        idx += 1

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    return output_path