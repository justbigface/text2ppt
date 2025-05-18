from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import os
import requests
from urllib.parse import urlparse

def get_img_path(img_url, default_img="default_cover.png"):
    """统一的图片路径处理函数
    
    Args:
        img_url: 图片URL或本地路径
        default_img: 默认图片路径，相对于项目根目录
        
    Returns:
        str: 有效的图片路径
    """
    if not img_url or not isinstance(img_url, str):
        return default_img
        
    # 如果是dummyimage.com的链接，直接使用默认图片
    if "dummyimage.com" in img_url.lower():
        return default_img
        
    # 如果是本地路径且存在，直接返回
    if os.path.exists(img_url):
        return img_url
        
    # 如果是http(s)链接，验证URL格式
    if img_url.startswith(('http://', 'https://')):
        try:
            parsed = urlparse(img_url)
            if parsed.scheme and parsed.netloc:
                return img_url
        except:
            pass
            
    return default_img

def add_placeholder_shape(slide, left, top, width, height, text="No Image"):
    """添加占位图形
    
    Args:
        slide: PPT幻灯片对象
        left: 左边距
        top: 上边距
        width: 宽度
        height: 高度
        text: 显示文本
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
    shape.line.color.rgb = RGBColor(200, 200, 200)
    
    # 添加文本
    text_frame = shape.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER