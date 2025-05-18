from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .utils import get_img_path, add_placeholder_shape
import os

def create_cover_big_image_ppt(title, subtitle="", img=None, output_path="output.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 大图
    if img:
        img_path = get_img_path(img)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
        else:
            add_placeholder_shape(slide, 0, 0, prs.slide_width, prs.slide_height, "No Cover Image")

    # 标题居中遮罩
    box_w = Inches(10)
    box_h = Inches(1.5)
    left = (prs.slide_width - box_w) / 2
    top = (prs.slide_height - box_h) / 2

    title_box = slide.shapes.add_textbox(left, top, box_w, box_h)
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 半透明黑色背景
    fill = title_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0,0,0)
    fill.transparency = 0.4

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(left, top + box_h, box_w, Inches(0.7))
        sub_frame = sub_box.text_frame
        sub_frame.text = subtitle
        p = sub_frame.paragraphs[0]
        p.font.size = Pt(26)
        p.font.color.rgb = RGBColor(240,240,240)
        p.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return output_path