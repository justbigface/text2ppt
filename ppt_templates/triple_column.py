from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from .utils import get_img_path, add_placeholder_shape

def create_triple_column_ppt(title, items, output_path="output.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 背景色可选
    # slide.background.fill.solid()
    # slide.background.fill.fore_color.rgb = RGBColor(245,245,245)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(4.5), Inches(1), Inches(4), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = 1  # Center

    # 三列坐标和宽高
    card_width = Inches(3.3)
    card_height = Inches(3.5)
    spacing = Inches(0.5)
    start_left = Inches(0.8)
    top_img = Inches(2.0)
    img_height = Inches(1.4)
    text_top = top_img + img_height + Inches(0.1)

    for idx in range(3):  # 只画三列
        if idx >= len(items):
            break
        item = items[idx]
        left = start_left + idx * (card_width + spacing)
        # 图片框
        img_url = item.get("img", "")
        img_path = get_img_path(img_url)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top_img, card_width, img_height)
        else:
            add_placeholder_shape(slide, left, top_img, card_width, img_height)
        # 小标题
        sub_box = slide.shapes.add_textbox(left, text_top, card_width, Inches(0.5))
        sub_frame = sub_box.text_frame
        sub_frame.text = item.get("subtitle", "")
        p = sub_frame.paragraphs[0]
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        # 描述
        desc_box = slide.shapes.add_textbox(left, text_top + Inches(0.5), card_width, Inches(1.2))
        desc_frame = desc_box.text_frame
        desc_frame.text = item.get("desc", "")
        p = desc_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(60, 60, 60)

    # 加下划线/圆点可选
    # slide.shapes.add_shape(...) 画线和圆点

    prs.save(output_path)
    return output_path