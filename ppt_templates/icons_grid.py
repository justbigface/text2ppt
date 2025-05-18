from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_icons_grid_ppt(title, items, output_path="output.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.5), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(38)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 计算宫格（最多3x3，建议6-9项）
    cols = 3
    rows = 3 if len(items) > 6 else 2
    grid_w = Inches(11)
    grid_h = Inches(5)
    left_margin = Inches(1.2)
    top_margin = Inches(1.8)
    cell_w = grid_w / cols
    cell_h = grid_h / rows

    for idx, item in enumerate(items[:cols*rows]):
        row = idx // cols
        col = idx % cols
        left = left_margin + col * cell_w
        top = top_margin + row * cell_h

        # 图标（如果是emoji直接显示，否则可用图片，默认圆圈图形）
        icon_text = item.get("icon", "●")
        icon_box = slide.shapes.add_textbox(left, top, cell_w, Inches(0.6))
        icon_frame = icon_box.text_frame
        icon_frame.text = icon_text
        p = icon_frame.paragraphs[0]
        p.font.size = Pt(38)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 标题
        sub_box = slide.shapes.add_textbox(left, top + Inches(0.7), cell_w, Inches(0.3))
        sub_frame = sub_box.text_frame
        sub_frame.text = item.get("subtitle", "")
        p = sub_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 描述
        desc_box = slide.shapes.add_textbox(left, top + Inches(1.1), cell_w, Inches(0.8))
        desc_frame = desc_box.text_frame
        desc_frame.text = item.get("desc", "")
        p = desc_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

    prs.save(output_path)
    return output_path