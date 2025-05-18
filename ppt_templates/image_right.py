from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from .utils import get_img_path, add_placeholder_shape
import os

def create_image_right_ppt(title, items, output_path="output.pptx", img_path="your_default_image.jpg"):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 顶部黑底标题
    title_box = slide.shapes.add_textbox(Inches(4.5), Inches(1), Inches(4), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p.alignment = PP_ALIGN.CENTER
    # 黑色底色
    fill = title_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0,0,0)

    # 右侧大图片
    img_path = get_img_path(img_path)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(8), Inches(2), Inches(4.5), Inches(3.7))
    else:
        add_placeholder_shape(slide, Inches(8), Inches(2), Inches(4.5), Inches(3.7))

    # 左侧内容区
    main_left = Inches(0.8)
    main_top = Inches(2.1)
    main_width = Inches(6.5)
    main_height = Inches(4.0)

    # 主要内容标题
    main_title_box = slide.shapes.add_textbox(main_left, main_top, main_width, Inches(0.7))
    main_title_frame = main_title_box.text_frame
    main_title_frame.word_wrap = True
    main_title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    main_title_frame.text = "主要内容"
    p = main_title_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255,255,255)
    p.alignment = PP_ALIGN.CENTER
    fill = main_title_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0,0,0)

    # 内容区，多段落（或带icon）
    item_top = main_top + Inches(0.8)
    for idx, item in enumerate(items):
        if idx > 3: break  # 最多4条
        # 图标省略，可用Unicode或用 add_picture
        box = slide.shapes.add_textbox(main_left, item_top + idx*Inches(0.8), main_width, Inches(0.7))
        frame = box.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.text = f"{item.get('subtitle','')}  {item.get('desc','')}"
        p = frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0,0,0)
        p.alignment = PP_ALIGN.LEFT

    prs.save(output_path)
    return output_path